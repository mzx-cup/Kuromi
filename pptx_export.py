"""
PPT导出模块（增强版）
将生成的课程数据导出为PPTX文件，支持多种元素类型
"""

import io
import os
import base64
import re
from typing import Optional, Any
from urllib.request import urlopen
from urllib.error import URLError

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

from state import CourseData, Slide, SlideElement, SlideBackground


class PPTXExporter:
    """
    PPTX导出器（增强版）
    将CourseData转换为PPTX文件，支持：
    - 文本、代码元素
    - 形状元素（矩形、圆形、三角形等）
    - 图表元素（柱状、饼图、折线图等）
    - 图片元素
    - 表格元素
    - LaTeX公式（通过SVG转换）
    """

    # 缩放比例：canvas 1000px -> PPTX 13.333 inches (96 DPI)
    PX_TO_INCH = 13.333 / 1000
    PX_TO_PT = 96 / 72  # 像素转磅值

    def __init__(self):
        self.prs = Presentation()
        # 设置16:9比例
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def export(self, course_data: CourseData) -> bytes:
        """导出课程数据为PPTX字节流"""
        # 添加标题页
        self._add_title_slide(course_data.title, course_data.teacher.name)

        # 添加每个幻灯片
        for slide in course_data.slides:
            self._add_slide(slide)

        # 添加结束页
        self._add_end_slide()

        # 保存到字节流
        pptx_stream = io.BytesIO()
        self.prs.save(pptx_stream)
        pptx_stream.seek(0)
        return pptx_stream.read()

    def export_to_file(self, course_data: CourseData, output_path: str) -> str:
        """导出课程数据到文件，返回文件路径"""
        pptx_bytes = self.export(course_data)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, 'wb') as f:
            f.write(pptx_bytes)
        return output_path

    def _px(self, value: float) -> Inches:
        """将像素值转换为英寸"""
        return Inches(value * self.PX_TO_INCH)

    def _parse_color(self, color_str: str) -> RgbColor:
        """解析颜色字符串为RgbColor"""
        if not color_str or color_str.startswith('#'):
            color_str = color_str.lstrip('#') if color_str else '333333'
            if len(color_str) == 6:
                r = int(color_str[0:2], 16)
                g = int(color_str[2:4], 16)
                b = int(color_str[4:6], 16)
                return RgbColor(r, g, b)
        return RgbColor(51, 51, 51)

    def _fetch_image_data(self, url: str) -> Optional[bytes]:
        """获取远程图片数据"""
        if not url:
            return None
        try:
            if url.startswith('data:image'):
                # Base64编码的图片
                match = re.match(r'data:image/[^;]+;base64,(.+)', url)
                if match:
                    return base64.b64decode(match.group(1))
            elif url.startswith(('http://', 'https://')):
                with urlopen(url, timeout=10) as response:
                    return response.read()
        except (URLError, Exception):
            pass
        return None

    def _add_title_slide(self, title: str, teacher_name: str):
        """添加标题页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(2.5),
            self.prs.slide_width, Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(114, 46, 209)
        shape.line.fill.background()

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.8),
            self.prs.slide_width - Inches(1), Inches(1)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(48)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 添加副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2),
            self.prs.slide_width - Inches(1), Inches(0.5)
        )
        tf = subtitle_box.text_frame
        tf.paragraphs[0].text = f"主讲教师: {teacher_name}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 添加Logo
        logo_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(2), Inches(0.5)
        )
        tf = logo_box.text_frame
        tf.paragraphs[0].text = "星识 Star-Learn"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = RgbColor(114, 46, 209)

    def _add_end_slide(self):
        """添加结束页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(114, 46, 209)
        shape.line.fill.background()

        thanks_box = slide.shapes.add_textbox(
            Inches(0), Inches(2.5),
            self.prs.slide_width, Inches(1)
        )
        tf = thanks_box.text_frame
        tf.paragraphs[0].text = "感谢观看"
        tf.paragraphs[0].font.size = Pt(56)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(
            Inches(0), Inches(4),
            self.prs.slide_width, Inches(0.5)
        )
        tf = subtitle_box.text_frame
        tf.paragraphs[0].text = "星识 Star-Learn - 全息智控学习舱"
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_slide(self, slide: Slide):
        """添加内容幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        pptx_slide = self.prs.slides.add_slide(slide_layout)

        # 设置背景
        self._set_slide_background(pptx_slide, slide.background)

        # 添加顶部装饰条
        header_shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(0.8)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RgbColor(114, 46, 209)
        header_shape.line.fill.background()

        # 添加幻灯片标题
        title_box = pptx_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(10), Inches(0.5)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = slide.title
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)

        # 添加页码
        page_box = pptx_slide.shapes.add_textbox(
            Inches(12.5), Inches(0.2),
            Inches(0.5), Inches(0.4)
        )
        tf = page_box.text_frame
        tf.paragraphs[0].text = str(slide.id)
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # 添加内容元素
        for elem in slide.content.elements:
            self._add_element(pptx_slide, elem)

    def _set_slide_background(self, slide, background: SlideBackground):
        """设置幻灯片背景"""
        if not background:
            return

        if background.type == "solid" and background.color:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = self._parse_color(background.color)
        elif background.type == "gradient" and background.gradient:
            # PPTX渐变背景简化处理为纯色
            colors = background.gradient.get("colors", [])
            if colors:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = self._parse_color(colors[0])

    def _add_element(self, slide, element: SlideElement):
        """根据元素类型添加到幻灯片"""
        elem_type = element.type

        if elem_type == "text":
            self._add_text_element(slide, element)
        elif elem_type == "code":
            self._add_code_element(slide, element)
        elif elem_type == "image":
            self._add_image_element(slide, element)
        elif elem_type == "shape":
            self._add_shape_element(slide, element)
        elif elem_type == "chart":
            self._add_chart_element(slide, element)
        elif elem_type == "table":
            self._add_table_element(slide, element)
        elif elem_type == "latex":
            self._add_latex_element(slide, element)
        elif elem_type == "line":
            self._add_line_element(slide, element)

    def _add_text_element(self, slide, element: SlideElement):
        """添加文本元素"""
        textbox = slide.shapes.add_textbox(
            self._px(element.left),
            self._px(element.top),
            self._px(element.width),
            self._px(element.height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True

        # 设置字体和颜色
        font_name = element.default_font_name or "Microsoft YaHei"
        font_color = self._parse_color(element.default_color or "#333333")

        # 处理多行文本
        lines = element.content.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = line.strip()
            p.font.size = Pt(18)
            p.font.color.rgb = font_color
            p.font.name = font_name
            p.space_after = Pt(8)

            # 处理列表符号
            if line.strip().startswith(('•', '○', '▸', '▫', '-', '*', '·')):
                p.level = 0

    def _add_code_element(self, slide, element: SlideElement):
        """添加代码元素"""
        # 添加深色背景
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self._px(element.left),
            self._px(element.top),
            self._px(element.width),
            self._px(element.height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(30, 30, 30)
        shape.line.color.rgb = RgbColor(100, 100, 100)

        # 添加代码文本
        code_box = slide.shapes.add_textbox(
            self._px(element.left) + Inches(0.1),
            self._px(element.top) + Inches(0.1),
            self._px(element.width) - Inches(0.2),
            self._px(element.height) - Inches(0.2)
        )
        tf = code_box.text_frame
        tf.word_wrap = True

        lines = element.content.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.font.name = "Monaco"
            p.font.color.rgb = RgbColor(212, 212, 212)

    def _add_image_element(self, slide, element: SlideElement):
        """添加图片元素"""
        if not element.image_url:
            return

        image_data = self._fetch_image_data(element.image_url)
        if not image_data:
            return

        try:
            slide.shapes.add_picture(
                io.BytesIO(image_data),
                self._px(element.left),
                self._px(element.top),
                self._px(element.width),
                self._px(element.height)
            )
        except Exception:
            pass

    def _add_shape_element(self, slide, element: SlideElement):
        """添加形状元素"""
        shape_name = element.shape_name.lower() if element.shape_name else "rectangle"

        # 映射形状名称
        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "circle": MSO_SHAPE.OVAL,
            "ellipse": MSO_SHAPE.ELLIPSE,
            "triangle": MSO_SHAPE.TRIANGLE,
            "pentagon": MSO_SHAPE.PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "star": MSO_SHAPE.STAR_5_POINT,
            "arrow": MSO_SHAPE.RIGHT_ARROW,
            "cylinder": MSO_SHAPE.CAN,
            "cube": MSO_SHAPE.CUBE,
        }

        ms_shape = shape_map.get(shape_name, MSO_SHAPE.RECTANGLE)

        try:
            shape = slide.shapes.add_shape(
                ms_shape,
                self._px(element.left),
                self._px(element.top),
                self._px(element.width),
                self._px(element.height)
            )

            # 设置填充色
            if element.fill:
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._parse_color(element.fill)
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RgbColor(91, 155, 213)

            # 设置边框
            if element.outline and element.outline.get("color"):
                shape.line.color.rgb = self._parse_color(element.outline.get("color", "#333333"))
                shape.line.width = Pt(element.outline.get("width", 1))
            else:
                shape.line.fill.background()

        except Exception:
            # 回退到矩形
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                self._px(element.left),
                self._px(element.top),
                self._px(element.width),
                self._px(element.height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RgbColor(91, 155, 213)
            shape.line.fill.background()

    def _add_chart_element(self, slide, element: SlideElement):
        """添加图表元素"""
        if not element.chart_data:
            return

        chart_type = element.chart_type.lower() if element.chart_type else "bar"
        chart_data = element.chart_data

        try:
            labels = chart_data.get("labels", [])
            series_data = chart_data.get("series", [])

            if not series_data:
                return

            # 创建图表数据
            cd = CategoryChartData()
            cd.categories = labels

            for i, series in enumerate(series_data):
                series_name = f"Series {i + 1}"
                cd.add_series(series_name, series)

            # 映射图表类型
            type_map = {
                "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "column": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
                "ring": XL_CHART_TYPE.DOUGHNUT,
                "area": XL_CHART_TYPE.AREA,
                "scatter": XL_CHART_TYPE.XY_SCATTER,
                "radar": XL_CHART_TYPE.RADAR,
            }

            xl_type = type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

            chart = slide.shapes.add_chart(
                xl_type,
                self._px(element.left),
                self._px(element.top),
                self._px(element.width),
                self._px(element.height),
                cd
            ).chart

            # 设置图例
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        except Exception as e:
            # 图表创建失败，添加占位文本
            textbox = slide.shapes.add_textbox(
                self._px(element.left),
                self._px(element.top),
                self._px(element.width),
                self._px(element.height)
            )
            tf = textbox.text_frame
            tf.paragraphs[0].text = f"[图表: {chart_type}]"
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = RgbColor(128, 128, 128)

    def _add_table_element(self, slide, element: SlideElement):
        """添加表格元素"""
        if not element.table_data:
            return

        rows = element.table_data
        if not rows:
            return

        num_rows = len(rows)
        num_cols = max(len(row) for row in rows) if rows else 0

        if num_rows == 0 or num_cols == 0:
            return

        try:
            table = slide.shapes.add_table(
                num_rows, num_cols,
                self._px(element.left),
                self._px(element.top),
                self._px(element.width),
                self._px(element.height)
            ).table

            # 设置列宽
            col_widths = element.col_widths or [1.0 / num_cols] * num_cols
            for i, width_ratio in enumerate(col_widths):
                if i < num_cols:
                    table.columns[i].width = int(self._px(element.width).emu * width_ratio)

            # 填充数据
            for row_idx, row in enumerate(rows):
                for col_idx, cell in enumerate(row):
                    if row_idx < num_rows and col_idx < num_cols:
                        cell_text = cell.get("text", "") if isinstance(cell, dict) else str(cell)
                        table.cell(row_idx, col_idx).text = cell_text

                        # 设置样式
                        para = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                        para.font.size = Pt(12)
                        para.font.name = "Microsoft YaHei"

                        # 表头样式
                        if row_idx == 0:
                            para.font.bold = True
                            para.font.color.rgb = RgbColor(255, 255, 255)
                            cell = table.cell(row_idx, col_idx)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RgbColor(99, 102, 241)

        except Exception:
            pass

    def _add_latex_element(self, slide, element: SlideElement):
        """添加LaTeX公式元素（作为文本显示）"""
        if not element.latex:
            return

        # LaTeX公式作为文本显示（PPTX原生不支持LaTeX）
        textbox = slide.shapes.add_textbox(
            self._px(element.left),
            self._px(element.top),
            self._px(element.width),
            self._px(element.height)
        )
        tf = textbox.text_frame
        tf.paragraphs[0].text = element.latex
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.name = "Cambria Math"
        tf.paragraphs[0].font.color.rgb = RgbColor(51, 51, 51)
        tf.word_wrap = True

    def _add_line_element(self, slide, element: SlideElement):
        """添加线条元素"""
        points = element.points
        if len(points) < 2:
            return

        try:
            start_x, start_y = points[0]
            end_x, end_y = points[1]

            connector = slide.shapes.add_connector(
                1,  # 直线连接器
                self._px(start_x),
                self._px(start_y),
                self._px(end_x),
                self._px(end_y)
            )

            connector.line.color.rgb = self._parse_color(element.line_color or "#333333")

            # 线条样式
            line_style = element.line_style.lower() if element.line_style else "solid"
            if line_style == "dashed":
                connector.line.dash_style = 2  # DASH
            elif line_style == "dotted":
                connector.line.dash_style = 3  # DOT

            connector.line.width = Pt(2)

        except Exception:
            pass


# 全局实例
_pptx_exporter: Optional[PPTXExporter] = None


def get_pptx_exporter() -> PPTXExporter:
    """获取PPTX导出器实例"""
    global _pptx_exporter
    if _pptx_exporter is None:
        _pptx_exporter = PPTXExporter()
    return _pptx_exporter


def export_course_to_pptx(course_data: CourseData) -> bytes:
    """快捷函数：将课程数据导出为PPTX"""
    exporter = get_pptx_exporter()
    return exporter.export(course_data)


def export_course_to_file(course_data: CourseData, output_path: str) -> str:
    """快捷函数：将课程数据导出到文件"""
    exporter = get_pptx_exporter()
    return exporter.export_to_file(course_data, output_path)