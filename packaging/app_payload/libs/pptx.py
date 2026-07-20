"""
PPT导出模块（增强版）
将生成的课程数据导出为PPTX文件，支持多种元素类型
"""

import io
import os
import base64
import re
from typing import Optional, Any
from urllib.request import urlopen
from urllib.error import URLError

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

from state import CourseData, Slide, SlideElement, SlideBackground, SlideV2, SlideContentItemV2


class PPTXExporter:
    """
    PPTX导出器（增强版）
    将CourseData转换为PPTX文件，支持：
    - 文本、代码元素
    - 形状元素（矩形、圆形、三角形等）
    - 图表元素（柱状、饼图、折线图等）
    - 图片元素
    - 表格元素
    - LaTeX公式（通过SVG转换）
    """

    # 缩放比例：canvas 1000px -> PPTX 13.333 inches (96 DPI)
    PX_TO_INCH = 13.333 / 1000
    PX_TO_PT = 96 / 72  # 像素转磅值

    def __init__(self):
        self.prs = Presentation()
        # 设置16:9比例
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def export(self, course_data: CourseData) -> bytes:
        """导出课程数据为PPTX字节流"""
        # 添加标题页
        self._add_title_slide(course_data.title, course_data.teacher.name)

        # 添加每个幻灯片
        for slide in course_data.slides:
            self._add_slide(slide)

        # 添加结束页
        self._add_end_slide()

        # 保存到字节流
        pptx_stream = io.BytesIO()
        self.prs.save(pptx_stream)
        pptx_stream.seek(0)
        return pptx_stream.read()

    def export_to_file(self, course_data: CourseData, output_path: str) -> str:
        """导出课程数据到文件，返回文件路径"""
        pptx_bytes = self.export(course_data)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, 'wb') as f:
            f.write(pptx_bytes)
        return output_path

    def export_v2(self, course_data: CourseData) -> bytes:
        """导出 V2 结构化布局的课程数据为 PPTX 字节流"""
        # 添加标题页
        self._add_title_slide(course_data.title, course_data.teacher.name)

        # 添加 V2 幻灯片
        for slide_v2 in course_data.slides_v2:
            self._add_slide_v2(slide_v2)

        # 添加传统幻灯片（如果有）
        for slide in course_data.slides:
            self._add_slide(slide)

        # 添加结束页
        self._add_end_slide()

        # 保存到字节流
        pptx_stream = io.BytesIO()
        self.prs.save(pptx_stream)
        pptx_stream.seek(0)
        return pptx_stream.read()

    def export_v2_to_file(self, course_data: CourseData, output_path: str) -> str:
        """导出 V2 课程数据到文件，返回文件路径"""
        pptx_bytes = self.export_v2(course_data)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, 'wb') as f:
            f.write(pptx_bytes)
        return output_path

    def _add_slide_v2(self, slide_v2: SlideV2):
        """添加 V2 结构化布局幻灯片 — 支持 OpenMAIC 和卡片两种格式, 20+ 种 layout"""
        # OpenMAIC 格式 (MiniMax PPT provider 生成): 有 elements 字段
        elements = getattr(slide_v2, 'elements', None) or []
        if elements:
            self._render_openmaic_slide(slide_v2)
            return

        slide_layout = self.prs.slide_layouts[6]
        pptx_slide = self.prs.slides.add_slide(slide_layout)

        # 使用 theme 颜色或默认配色
        theme = getattr(slide_v2, 'theme', None) or {}
        bg_info = getattr(slide_v2, 'background', None) or {}
        bg_color = self._hex_to_rgb(bg_info.get('color', '#FFFFFF')) if bg_info.get('type') == 'solid' else RGBColor(255, 255, 255)
        header_color = self._hex_to_rgb(theme.get('primary', '#1E40AF'))
        accent_color = self._hex_to_rgb(theme.get('accent', '#3B82F6'))
        text_on_header = RGBColor(255, 255, 255)

        # 背景
        pptx_slide.background.fill.solid()
        pptx_slide.background.fill.fore_color.rgb = bg_color

        # 根据 layoutType 计算卡片布局
        layout_type = slide_v2.layout_type or 'two-column'
        cards = slide_v2.content or []

        # title-only / hero-center / edu-welcome: 不画 header, 大标题居中
        if layout_type in ('title-only', 'hero-center', 'edu-welcome', 'spotlight-focus'):
            self._add_v2_layout_title_centered(pptx_slide, slide_v2.title, cards, accent_color, theme)
        elif layout_type in ('kinetic-type',):
            self._add_v2_layout_kinetic(pptx_slide, slide_v2.title, cards, accent_color)
        elif layout_type in ('orbit-ring', 'circle-radial'):
            self._add_v2_layout_radial(pptx_slide, slide_v2.title, cards, accent_color)
        elif layout_type in ('stair-step',):
            self._add_v2_layout_stair(pptx_slide, slide_v2.title, cards, accent_color)
        elif layout_type in ('quote-wall', 'quote-highlight'):
            self._add_v2_layout_quote(pptx_slide, slide_v2.title, cards, accent_color)
        elif layout_type in ('info-graphic',):
            self._add_v2_layout_infographic(pptx_slide, slide_v2.title, cards, accent_color)
        elif layout_type in ('isometric-cards', 'floating-overlap'):
            self._add_v2_layout_iso_cards(pptx_slide, slide_v2.title, cards, accent_color)
        elif layout_type in ('gradient-split',):
            self._add_v2_layout_gradient_split(pptx_slide, slide_v2.title, cards, accent_color, header_color)
        elif layout_type in ('dark-header',):
            self._add_v2_layout_dark_header(pptx_slide, slide_v2.title, cards, accent_color, text_on_header)
        elif layout_type in ('edu-definition',):
            self._add_v2_layout_edu_definition(pptx_slide, slide_v2.title, cards, accent_color)
        elif layout_type in ('edu-example',):
            self._add_v2_layout_edu_example(pptx_slide, slide_v2.title, cards, accent_color)
        else:
            # 通用 header 模式: 含 header bar + 卡片区
            self._add_v2_layout_with_header(pptx_slide, slide_v2.title, cards, layout_type, header_color, accent_color, text_on_header, theme)

    # ------------------------------------------------------------------
    # Layout-specific renderers — 每个 layout 一个差异化实现
    # ------------------------------------------------------------------

    def _add_v2_layout_title_centered(self, slide, title, cards, accent, theme):
        """title-only / hero-center / edu-welcome / spotlight-focus: 大标题居中"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title or ""
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = self._hex_to_rgb(theme.get('text', '#1E293B'))

        if cards:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.33), Inches(2.5))
            stf = sub_box.text_frame
            stf.word_wrap = True
            for i, card in enumerate(cards[:3]):
                p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = card.sub_title or card.text or ""
                run.font.size = Pt(20)
                run.font.color.rgb = self._hex_to_rgb(theme.get('text_secondary', '#64748B'))

    def _add_v2_layout_kinetic(self, slide, title, cards, accent):
        """kinetic-type: 倾斜大标题 + 装饰条"""
        # 装饰斜条
        deco = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(0.5), Inches(1.2), Inches(2), Inches(0.3))
        deco.fill.solid()
        deco.fill.fore_color.rgb = accent
        deco.line.fill.background()
        # 大标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.33), Inches(1.2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title or ""
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.italic = True
        run.font.color.rgb = RGBColor(15, 23, 42)
        # 卡片水平排列
        for i, card in enumerate(cards[:4]):
            left = Inches(0.5 + i * 3.1)
            self._add_v2_card(slide, card, left, Inches(3.5), Inches(2.9), Inches(2.5))

    def _add_v2_layout_radial(self, slide, title, cards, accent):
        """orbit-ring / circle-radial: 中央大圆 + 6 个小圆环绕"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
        p = title_box.text_frame.paragraphs[0]
        p.text = title or ""
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(15, 23, 42)
        # 中心圆
        center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(2.5), Inches(3.33), Inches(2.0))
        center.fill.solid()
        center.fill.fore_color.rgb = accent
        center.line.fill.background()
        ctitle = slide.shapes.add_textbox(Inches(5.5), Inches(3.1), Inches(3.33), Inches(0.8))
        ctf = ctitle.text_frame
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cp.text = title or "中心"
        cp.font.size = Pt(20)
        cp.font.bold = True
        cp.font.color.rgb = RGBColor(255, 255, 255)
        # 6 个环绕小圆 (12, 2, 4, 6, 8, 10 点钟方向)
        import math
        for i, card in enumerate(cards[:6]):
            angle = (i * 60 - 90) * math.pi / 180
            cx = 7.16 + 4.5 * math.cos(angle)
            cy = 3.5 + 1.8 * math.sin(angle)
            sz = 1.4
            sat = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - sz/2), Inches(cy - sz/2), Inches(sz), Inches(sz))
            sat.fill.solid()
            sat.fill.fore_color.rgb = RGBColor(255, 255, 255)
            sat.line.color.rgb = accent
            sat.line.width = Pt(2)
            tbox = slide.shapes.add_textbox(Inches(cx - sz/2), Inches(cy - 0.3), Inches(sz), Inches(0.6))
            tp = tbox.text_frame.paragraphs[0]
            tp.alignment = PP_ALIGN.CENTER
            tp.text = (card.sub_title or card.text or "")[:8]
            tp.font.size = Pt(10)
            tp.font.color.rgb = RGBColor(15, 23, 42)

    def _add_v2_layout_stair(self, slide, title, cards, accent):
        """stair-step: 阶梯式排列 (4 级台阶)"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
        p = title_box.text_frame.paragraphs[0]
        p.text = title or ""
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(15, 23, 42)
        for i, card in enumerate(cards[:4]):
            level = i
            top = 1.2 + level * 0.9
            left = 0.5 + level * 1.5
            w = 12.33 - level * 1.5
            self._add_v2_card(slide, card, Inches(left), Inches(top), Inches(w), Inches(0.85))

    def _add_v2_layout_quote(self, slide, title, cards, accent):
        """quote-wall / quote-highlight: 大引号 + 引用文本"""
        # 大引号
        quote_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2), Inches(2))
        qp = quote_box.text_frame.paragraphs[0]
        qp.text = "“"
        qp.font.size = Pt(120)
        qp.font.color.rgb = accent
        # 引用文本
        if cards:
            ref_box = slide.shapes.add_textbox(Inches(2.5), Inches(1.8), Inches(10), Inches(2.5))
            rtf = ref_box.text_frame
            rtf.word_wrap = True
            rp = rtf.paragraphs[0]
            rp.text = cards[0].sub_title or cards[0].text or title or ""
            rp.font.size = Pt(28)
            rp.font.italic = True
            rp.font.color.rgb = RGBColor(15, 23, 42)
        # 标题
        if title and len(cards) > 1:
            tbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.33), Inches(0.5))
            tp = tbox.text_frame.paragraphs[0]
            tp.alignment = PP_ALIGN.RIGHT
            tp.text = f"— {title}"
            tp.font.size = Pt(14)
            tp.font.color.rgb = RGBColor(100, 116, 139)

    def _add_v2_layout_infographic(self, slide, title, cards, accent):
        """info-graphic: 大数字 + 描述"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
        p = title_box.text_frame.paragraphs[0]
        p.text = title or ""
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(71, 85, 105)
        # 3 大数字块
        for i, card in enumerate(cards[:3]):
            left = 0.5 + i * 4.2
            block = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.5), Inches(3.9), Inches(3.5))
            block.fill.solid()
            block.fill.fore_color.rgb = RGBColor(248, 250, 252)
            block.line.color.rgb = accent
            block.line.width = Pt(2)
            # 大数字
            num_box = slide.shapes.add_textbox(Inches(left), Inches(1.8), Inches(3.9), Inches(1.5))
            np_ = num_box.text_frame.paragraphs[0]
            np_.alignment = PP_ALIGN.CENTER
            run = np_.add_run()
            run.text = f"0{i+1}" if i < 9 else f"{i+1}"
            run.font.size = Pt(72)
            run.font.bold = True
            run.font.color.rgb = accent
            # 描述
            desc_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(3.5), Inches(3.5), Inches(1.3))
            dp = desc_box.text_frame.paragraphs[0]
            dp.alignment = PP_ALIGN.CENTER
            dp.word_wrap = True
            dp.text = (card.sub_title or card.text or "")[:50]
            dp.font.size = Pt(13)
            dp.font.color.rgb = RGBColor(71, 85, 105)

    def _add_v2_layout_iso_cards(self, slide, title, cards, accent):
        """isometric-cards / floating-overlap: 卡片层叠"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
        p = title_box.text_frame.paragraphs[0]
        p.text = title or ""
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(15, 23, 42)
        # 重叠卡片 (向右上偏移)
        for i, card in enumerate(cards[:4]):
            offset = i * 0.3
            self._add_v2_card(slide, card, Inches(0.7 + offset), Inches(1.5 + offset), Inches(8), Inches(1.2))

    def _add_v2_layout_gradient_split(self, slide, title, cards, accent, header):
        """gradient-split: 左右渐变分割"""
        # 左半深色块
        left_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(6.16), Inches(7.5))
        left_block.fill.solid()
        left_block.fill.fore_color.rgb = header
        left_block.line.fill.background()
        # 标题 (白字)
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(5.16), Inches(2.5))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title or ""
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        # 右侧卡片
        for i, card in enumerate(cards[:3]):
            self._add_v2_card(slide, card, Inches(6.5), Inches(1.2 + i * 1.8), Inches(6.3), Inches(1.5))

    def _add_v2_layout_dark_header(self, slide, title, cards, accent, text_on_header):
        """dark-header: 黑色 header bar 占 1/3, 下方内容"""
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), self.prs.slide_width, Inches(2.5))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
        header_shape.line.fill.background()
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(1))
        p = tbox.text_frame.paragraphs[0]
        p.text = title or ""
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = text_on_header
        # 下方卡片
        for i, card in enumerate(cards[:3]):
            self._add_v2_card(slide, card, Inches(0.5), Inches(2.9 + i * 1.5), Inches(12.33), Inches(1.3))

    def _add_v2_layout_edu_definition(self, slide, title, cards, accent):
        """edu-definition: 左侧定义框 + 右侧属性标签"""
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
        p = tbox.text_frame.paragraphs[0]
        p.text = title or ""
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(15, 23, 42)
        if not cards:
            return
        # 左定义框
        left_def = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5))
        left_def.fill.solid()
        left_def.fill.fore_color.rgb = RGBColor(219, 234, 254)
        left_def.line.fill.background()
        ltxt = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.6), Inches(5.1))
        ltf = ltxt.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.text = cards[0].sub_title or cards[0].text or ""
        lp.font.size = Pt(18)
        lp.font.color.rgb = RGBColor(30, 64, 175)
        # 右侧属性标签
        for i, card in enumerate(cards[1:4]):
            tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.4 + i * 1.7), Inches(5.83), Inches(1.4))
            tag.fill.solid()
            tag.fill.fore_color.rgb = RGBColor(254, 243, 199)
            tag.line.color.rgb = RGBColor(245, 158, 11)
            tag.line.width = Pt(1.5)
            tbox2 = slide.shapes.add_textbox(Inches(7.2), Inches(1.5 + i * 1.7), Inches(5.43), Inches(1.2))
            ttf = tbox2.text_frame
            ttf.word_wrap = True
            tp = ttf.paragraphs[0]
            tp.text = (card.sub_title or card.text or "")[:80]
            tp.font.size = Pt(13)
            tp.font.color.rgb = RGBColor(120, 53, 15)

    def _add_v2_layout_edu_example(self, slide, title, cards, accent):
        """edu-example: 左侧概念 + 右侧示例区"""
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
        p = tbox.text_frame.paragraphs[0]
        p.text = title or ""
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(15, 23, 42)
        if not cards:
            return
        # 左概念
        lt = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
        ltf = lt.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.text = cards[0].sub_title or cards[0].text or ""
        lp.font.size = Pt(16)
        lp.font.color.rgb = RGBColor(30, 41, 59)
        # 右示例
        for i, card in enumerate(cards[1:3]):
            ex = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.2 + i * 2.6), Inches(6), Inches(2.3))
            ex.fill.solid()
            ex.fill.fore_color.rgb = RGBColor(13, 17, 23)
            ex.line.fill.background()
            et = slide.shapes.add_textbox(Inches(7.0), Inches(1.4 + i * 2.6), Inches(5.6), Inches(2.0))
            etf = et.text_frame
            etf.word_wrap = True
            ep = etf.paragraphs[0]
            ep.text = (card.text or card.sub_title or "")[:200]
            ep.font.size = Pt(12)
            ep.font.name = 'Consolas'
            ep.font.color.rgb = RGBColor(201, 209, 217)

    def _add_v2_layout_with_header(self, slide, title, cards, layout_type, header_color, accent_color, text_on_header, theme):
        """通用 layout: header bar + 卡片区, 按 layout_type 选择内部排布"""
        # 标题栏
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), self.prs.slide_width, Inches(0.8))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = header_color
        header_shape.line.fill.background()
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.5))
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = text_on_header

        # 内容区按 layout 分支
        if layout_type in ('two-column', 'comparison', 'edu-keypoints'):
            self._add_v2_cards_two_column(slide, cards)
        elif layout_type in ('grid-cards',):
            self._add_v2_cards_grid(slide, cards)
        elif layout_type in ('code-showcase', 'terminal-style', 'concept-code', 'api-doc',
                             'edu-summary', 'edu-programming-concept'):
            # 代码/总结类: 第一卡片做代码块, 后续做要点
            self._add_v2_cards_vertical(slide, cards)
        else:
            # header-content / numbered-list / chapter-divider / step-by-step / etc.
            self._add_v2_cards_vertical(slide, cards)

    def _render_openmaic_slide(self, slide_v2: SlideV2):
        """渲染 OpenMAIC 格式的幻灯片 (MiniMax PPT provider 产出)"""
        slide_layout = self.prs.slide_layouts[6]
        pptx_slide = self.prs.slides.add_slide(slide_layout)

        # 背景
        bg_info = getattr(slide_v2, 'background', None) or {}
        if bg_info.get('type') == 'solid':
            bg_color = self._hex_to_rgb(bg_info.get('color', '#0F172A'))
        else:
            bg_color = RGBColor(15, 23, 42)
        pptx_slide.background.fill.solid()
        pptx_slide.background.fill.fore_color.rgb = bg_color

        # 主题: 用后端传入的 theme 决定标题色, 不再 hardcoded 浅灰
        theme = getattr(slide_v2, 'theme', None) or {}
        theme_name = theme.get('name', '') if isinstance(theme, dict) else ''
        if not theme_name and isinstance(theme, str):
            theme_name = theme
        # 探测是否是深色背景: 亮度 < 128 视为深色, 标题用浅色
        bg_luminance = (bg_color[0] * 299 + bg_color[1] * 587 + bg_color[2] * 114) / 1000
        if bg_luminance < 128:
            title_color = RGBColor(248, 250, 252)  # near-white for dark bg
            body_color = RGBColor(226, 232, 240)
        else:
            title_color = RGBColor(15, 23, 42)  # near-black for light bg
            body_color = RGBColor(30, 41, 59)

        # 标题
        title = getattr(slide_v2, 'title', '') or ''
        if title:
            title_box = pptx_slide.shapes.add_textbox(
                Inches(0.5), Inches(0.2),
                Inches(12), Inches(0.6)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = title_color

        # 遍历 elements 数组按序渲染
        elements = getattr(slide_v2, 'elements', None) or []
        # MiniMax 视口: 960 x 540 (16:9)
        VP_W, VP_H = 960.0, 540.0
        sw = self.prs.slide_width  # in EMU
        sh = self.prs.slide_height
        scale_x = sw / VP_W
        scale_y = sh / VP_H

        for el in elements:
            el_type = el.get('type', '')
            left = int(el.get('left', 0) * scale_x)
            top = int(el.get('top', 0) * scale_y)
            w = int(el.get('width', 100) * scale_x)
            h = int(el.get('height', 40) * scale_y)
            fill_color = el.get('fill', 'transparent')

            if el_type == 'shape':
                shape_name = el.get('shape_name', 'rectangle')
                shape = pptx_slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE if shape_name == 'round-rectangle' else MSO_SHAPE.RECTANGLE,
                    left, top, w, h
                )
                if fill_color and fill_color != 'transparent':
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = self._hex_to_rgb(fill_color)
                else:
                    shape.fill.background()
                shape.line.fill.background()
                opacity = el.get('opacity', 1.0)
                if opacity < 1.0:
                    try:
                        shape.fill.fore_color.brightness = opacity
                    except Exception:
                        pass

            elif el_type == 'text':
                content = el.get('content', '')
                # Strip HTML tags for PPTX
                plain = self._strip_html(content)
                dfont = el.get('defaultFontName', 'Microsoft YaHei')
                # 优先用 element 显式 color, 否则用 theme 推导的 body_color, 再 fallback
                dcolor = el.get('defaultColor') or ('#' + ''.join(f'{c:02X}' for c in body_color))
                text_box = pptx_slide.shapes.add_textbox(left, top, w, h)
                tf = text_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = plain[:500]
                p.font.size = Pt(12)
                p.font.name = dfont
                p.font.color.rgb = self._hex_to_rgb(dcolor)

            elif el_type == 'code':
                code_text = el.get('content', '')
                # 代码块背景
                code_bg = pptx_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
                code_bg.fill.solid()
                code_bg.fill.fore_color.rgb = RGBColor(13, 17, 23)
                code_bg.line.fill.background()
                # 代码文本
                code_box = pptx_slide.shapes.add_textbox(
                    left + Inches(0.1), top + Inches(0.05),
                    w - Inches(0.2), h - Inches(0.1)
                )
                tf = code_box.text_frame
                tf.word_wrap = True
                lines = code_text.split('\n')[:30]
                for i, line in enumerate(lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line[:120]
                    p.font.name = 'Consolas'
                    p.font.size = Pt(9)
                    p.font.color.rgb = RGBColor(201, 209, 217)

            elif el_type == 'image':
                src_url = el.get('src', '')
                if src_url:
                    img_data = self._fetch_image_data(src_url)
                    if img_data:
                        try:
                            pptx_slide.shapes.add_picture(io.BytesIO(img_data), left, top, w, h)
                        except Exception:
                            pass

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """将 #RRGGBB 或 #RGB 转换为 RGBColor, 默认返回白色"""
        h = hex_color.lstrip('#')
        if len(h) == 3:
            h = ''.join(c * 2 for c in h)
        try:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except (ValueError, IndexError):
            return RGBColor(255, 255, 255)

    def _strip_html(self, html: str) -> str:
        """剥离 HTML 标签, 保留纯文本"""
        import re
        # 替换 <br> 为换行
        text = re.sub(r'<br\s*/?>', '\n', html, flags=re.IGNORECASE)
        # 移除所有标签
        text = re.sub(r'<[^>]+>', '', text)
        # 解码常见实体
        text = text.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
        text = text.replace('&quot;', '"').replace('&#39;', "'").replace('&nbsp;', ' ')
        return text

    def _add_v2_cards_two_column(self, slide, cards):
        """添加两栏布局的卡片"""
        col_width = 5.5
        row_height = 2.5
        start_y = 1.2

        for idx, card in enumerate(cards[:2]):  # 最多2列
            left = Inches(0.5) if idx == 0 else Inches(6.8)
            self._add_v2_card(slide, card, left, Inches(start_y), Inches(col_width), Inches(row_height))

    def _add_v2_cards_grid(self, slide, cards):
        """添加网格布局的卡片"""
        col_width = 5.5
        row_height = 2.5
        start_x = [0.5, 6.8]
        start_y = [1.2, 3.9]

        for idx, card in enumerate(cards[:4]):  # 最多2x2网格
            col = idx % 2
            row = idx // 2
            left = Inches(start_x[col])
            top = Inches(start_y[row])
            self._add_v2_card(slide, card, left, top, Inches(col_width), Inches(row_height))

    def _add_v2_cards_vertical(self, slide, cards):
        """添加垂直布局的卡片"""
        col_width = 11.5
        row_height = 1.8
        start_y = 1.2

        for idx, card in enumerate(cards[:3]):  # 最多3个
            self._add_v2_card(slide, card, Inches(0.5), Inches(start_y + idx * 2), Inches(col_width), Inches(row_height))

    def _add_v2_card(self, slide, card: SlideContentItemV2, left, top, width, height):
        """添加单个 V2 内容卡片"""
        # 卡片背景
        bg_color = self._parse_v2_color(card.color_theme or 'blue')
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = bg_color
        bg_shape.line.fill.background()

        # 内边距区域
        text_left = left + Inches(0.15)
        text_top = top + Inches(0.1)
        text_width = width - Inches(0.3)
        text_height = height - Inches(0.2)

        # 添加小标题
        if card.sub_title:
            title_box = slide.shapes.add_textbox(text_left, text_top, text_width, Inches(0.4))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"📖 {card.sub_title}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self._get_text_color(card.color_theme or 'blue')

        # 添加正文
        if card.text:
            text_box = slide.shapes.add_textbox(
                text_left,
                text_top + Inches(0.4),
                text_width,
                Inches(height.pt - Inches(0.5).pt)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = card.text[:200]  # 限制长度
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(30, 41, 59)  # #1E293B

        # 添加代码块
        if card.code_snippet:
            code_box = slide.shapes.add_textbox(
                text_left,
                text_top + Inches(height.pt - Inches(0.8).pt),
                text_width,
                Inches(0.7)
            )
            tf = code_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = card.code_snippet[:100]
            p.font.name = 'Consolas'
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(255, 255, 255)

        # 添加配图
        if card.image_url:
            img_data = self._fetch_image_data(card.image_url)
            if img_data:
                try:
                    pic = slide.shapes.add_picture(
                        io.BytesIO(img_data),
                        left + Inches(0.1),
                        top + Inches(height.pt - Inches(1.2).pt),
                        width - Inches(0.2),
                        Inches(1)
                    )
                except Exception:
                    pass

        # 添加视频（仅记录视频URL作为文本标签）
        if card.video_url:
            video_label_box = slide.shapes.add_textbox(
                text_left,
                text_top + Inches(0.5),
                text_width,
                Inches(0.4)
            )
            tf = video_label_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"🎬 视频: {card.video_url[:60]}..."
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(100, 116, 139)
            p.font.italic = True

    def _parse_v2_color(self, color_theme: str) -> RGBColor:
        """将 V2 色系转换为 RGBColor"""
        color_map = {
            'blue': (0xDB, 0xEA, 0xFE),      # #DBEAFE
            'yellow': (0xFE, 0xF3, 0xC7),     # #FEF3C7
            'green': (0xD1, 0xFA, 0xE5),      # #D1FAE5
            'purple': (0xED, 0xE9, 0xFE),     # #EDE9FE
            'orange': (0xFF, 0xED, 0xD5),      # #FFEDD5
        }
        rgb = color_map.get(color_theme, (0xDB, 0xEA, 0xFE))
        return RGBColor(*rgb)

    def _get_text_color(self, color_theme: str) -> RGBColor:
        """获取 V2 色系对应的文字颜色"""
        color_map = {
            'blue': RGBColor(30, 64, 175),    # #1E40AF
            'yellow': RGBColor(146, 64, 14),  # #92400E
            'green': RGBColor(6, 95, 70),     # #065F46
            'purple': RGBColor(109, 40, 217),  # #6D28D9
            'orange': RGBColor(194, 65, 12),   # #C2410C
        }
        return color_map.get(color_theme, RGBColor(30, 64, 175))

    def _px(self, value: float) -> Inches:
        """将像素值转换为英寸"""
        return Inches(value * self.PX_TO_INCH)

    def _parse_color(self, color_str: str) -> RGBColor:
        """解析颜色字符串为RGBColor"""
        if not color_str or color_str.startswith('#'):
            color_str = color_str.lstrip('#') if color_str else '333333'
            if len(color_str) == 6:
                r = int(color_str[0:2], 16)
                g = int(color_str[2:4], 16)
                b = int(color_str[4:6], 16)
                return RGBColor(r, g, b)
        return RGBColor(51, 51, 51)

    def _fetch_image_data(self, url: str) -> Optional[bytes]:
        """获取远程图片数据"""
        if not url:
            return None
        try:
            if url.startswith('data:image'):
                # Base64编码的图片
                match = re.match(r'data:image/[^;]+;base64,(.+)', url)
                if match:
                    return base64.b64decode(match.group(1))
            elif url.startswith(('http://', 'https://')):
                with urlopen(url, timeout=10) as response:
                    return response.read()
        except (URLError, Exception):
            pass
        return None

    def _add_title_slide(self, title: str, teacher_name: str):
        """添加标题页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(2.5),
            self.prs.slide_width, Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(114, 46, 209)
        shape.line.fill.background()

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.8),
            self.prs.slide_width - Inches(1), Inches(1)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(48)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 添加副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2),
            self.prs.slide_width - Inches(1), Inches(0.5)
        )
        tf = subtitle_box.text_frame
        tf.paragraphs[0].text = f"主讲教师: {teacher_name}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 添加Logo
        logo_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(2), Inches(0.5)
        )
        tf = logo_box.text_frame
        tf.paragraphs[0].text = "星识 Star-Learn"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = RGBColor(114, 46, 209)

    def _add_end_slide(self):
        """添加结束页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(114, 46, 209)
        shape.line.fill.background()

        thanks_box = slide.shapes.add_textbox(
            Inches(0), Inches(2.5),
            self.prs.slide_width, Inches(1)
        )
        tf = thanks_box.text_frame
        tf.paragraphs[0].text = "感谢观看"
        tf.paragraphs[0].font.size = Pt(56)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(
            Inches(0), Inches(4),
            self.prs.slide_width, Inches(0.5)
        )
        tf = subtitle_box.text_frame
        tf.paragraphs[0].text = "星识 Star-Learn - 全息智控学习舱"
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_slide(self, slide: Slide):
        """添加内容幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        pptx_slide = self.prs.slides.add_slide(slide_layout)

        # 设置背景
        self._set_slide_background(pptx_slide, slide.background)

        # 添加顶部装饰条
        header_shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(0.8)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(114, 46, 209)
        header_shape.line.fill.background()

        # 添加幻灯片标题
        title_box = pptx_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(10), Inches(0.5)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = slide.title
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # 添加页码
        page_box = pptx_slide.shapes.add_textbox(
            Inches(12.5), Inches(0.2),
            Inches(0.5), Inches(0.4)
        )
        tf = page_box.text_frame
        tf.paragraphs[0].text = str(slide.id)
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # 添加内容元素
        for elem in slide.content.elements:
            self._add_element(pptx_slide, elem)

    def _set_slide_background(self, slide, background: SlideBackground):
        """设置幻灯片背景"""
        if not background:
            return

        if background.type == "solid" and background.color:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = self._parse_color(background.color)
        elif background.type == "gradient" and background.gradient:
            # PPTX渐变背景简化处理为纯色
            colors = background.gradient.get("colors", [])
            if colors:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = self._parse_color(colors[0])

    def _add_element(self, slide, element: SlideElement):
        """根据元素类型添加到幻灯片"""
        elem_type = element.type

        if elem_type == "text":
            self._add_text_element(slide, element)
        elif elem_type == "code":
            self._add_code_element(slide, element)
        elif elem_type == "image":
            self._add_image_element(slide, element)
        elif elem_type == "shape":
            self._add_shape_element(slide, element)
        elif elem_type == "chart":
            self._add_chart_element(slide, element)
        elif elem_type == "table":
            self._add_table_element(slide, element)
        elif elem_type == "latex":
            self._add_latex_element(slide, element)
        elif elem_type == "line":
            self._add_line_element(slide, element)

    def _add_text_element(self, slide, element: SlideElement):
        """添加文本元素"""
        textbox = slide.shapes.add_textbox(
            self._px(element.left),
            self._px(element.top),
            self._px(element.width),
            self._px(element.height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True

        # 设置字体和颜色
        font_name = element.default_font_name or "Microsoft YaHei"
        font_color = self._parse_color(element.default_color or "#333333")

        # 处理多行文本
        lines = element.content.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = line.strip()
            p.font.size = Pt(18)
            p.font.color.rgb = font_color
            p.font.name = font_name
            p.space_after = Pt(8)

            # 处理列表符号
            if line.strip().startswith(('•', '○', '▸', '▫', '-', '*', '·')):
                p.level = 0

    def _add_code_element(self, slide, element: SlideElement):
        """添加代码元素"""
        # 添加深色背景
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self._px(element.left),
            self._px(element.top),
            self._px(element.width),
            self._px(element.height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 30, 30)
        shape.line.color.rgb = RGBColor(100, 100, 100)

        # 添加代码文本
        code_box = slide.shapes.add_textbox(
            self._px(element.left) + Inches(0.1),
            self._px(element.top) + Inches(0.1),
            self._px(element.width) - Inches(0.2),
            self._px(element.height) - Inches(0.2)
        )
        tf = code_box.text_frame
        tf.word_wrap = True

        lines = element.content.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.font.name = "Monaco"
            p.font.color.rgb = RGBColor(212, 212, 212)

    def _add_image_element(self, slide, element: SlideElement):
        """添加图片元素"""
        if not element.image_url:
            return

        image_data = self._fetch_image_data(element.image_url)
        if not image_data:
            return

        try:
            slide.shapes.add_picture(
                io.BytesIO(image_data),
                self._px(element.left),
                self._px(element.top),
                self._px(element.width),
                self._px(element.height)
            )
        except Exception:
            pass

    def _add_shape_element(self, slide, element: SlideElement):
        """添加形状元素"""
        shape_name = element.shape_name.lower() if element.shape_name else "rectangle"

        # 映射形状名称
        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "circle": MSO_SHAPE.OVAL,
            "ellipse": MSO_SHAPE.ELLIPSE,
            "triangle": MSO_SHAPE.TRIANGLE,
            "pentagon": MSO_SHAPE.PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "star": MSO_SHAPE.STAR_5_POINT,
            "arrow": MSO_SHAPE.RIGHT_ARROW,
            "cylinder": MSO_SHAPE.CAN,
            "cube": MSO_SHAPE.CUBE,
        }

        ms_shape = shape_map.get(shape_name, MSO_SHAPE.RECTANGLE)

        try:
            shape = slide.shapes.add_shape(
                ms_shape,
                self._px(element.left),
                self._px(element.top),
                self._px(element.width),
                self._px(element.height)
            )

            # 设置填充色
            if element.fill:
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._parse_color(element.fill)
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(91, 155, 213)

            # 设置边框
            if element.outline and element.outline.get("color"):
                shape.line.color.rgb = self._parse_color(element.outline.get("color", "#333333"))
                shape.line.width = Pt(element.outline.get("width", 1))
            else:
                shape.line.fill.background()

        except Exception:
            # 回退到矩形
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                self._px(element.left),
                self._px(element.top),
                self._px(element.width),
                self._px(element.height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(91, 155, 213)
            shape.line.fill.background()

    def _add_chart_element(self, slide, element: SlideElement):
        """添加图表元素"""
        if not element.chart_data:
            return

        chart_type = element.chart_type.lower() if element.chart_type else "bar"
        chart_data = element.chart_data

        try:
            labels = chart_data.get("labels", [])
            series_data = chart_data.get("series", [])

            if not series_data:
                return

            # 创建图表数据
            cd = CategoryChartData()
            cd.categories = labels

            for i, series in enumerate(series_data):
                series_name = f"Series {i + 1}"
                cd.add_series(series_name, series)

            # 映射图表类型
            type_map = {
                "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "column": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
                "ring": XL_CHART_TYPE.DOUGHNUT,
                "area": XL_CHART_TYPE.AREA,
                "scatter": XL_CHART_TYPE.XY_SCATTER,
                "radar": XL_CHART_TYPE.RADAR,
            }

            xl_type = type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

            chart = slide.shapes.add_chart(
                xl_type,
                self._px(element.left),
                self._px(element.top),
                self._px(element.width),
                self._px(element.height),
                cd
            ).chart

            # 设置图例
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        except Exception as e:
            # 图表创建失败，添加占位文本
            textbox = slide.shapes.add_textbox(
                self._px(element.left),
                self._px(element.top),
                self._px(element.width),
                self._px(element.height)
            )
            tf = textbox.text_frame
            tf.paragraphs[0].text = f"[图表: {chart_type}]"
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

    def _add_table_element(self, slide, element: SlideElement):
        """添加表格元素"""
        if not element.table_data:
            return

        rows = element.table_data
        if not rows:
            return

        num_rows = len(rows)
        num_cols = max(len(row) for row in rows) if rows else 0

        if num_rows == 0 or num_cols == 0:
            return

        try:
            table = slide.shapes.add_table(
                num_rows, num_cols,
                self._px(element.left),
                self._px(element.top),
                self._px(element.width),
                self._px(element.height)
            ).table

            # 设置列宽
            col_widths = element.col_widths or [1.0 / num_cols] * num_cols
            for i, width_ratio in enumerate(col_widths):
                if i < num_cols:
                    table.columns[i].width = int(self._px(element.width).emu * width_ratio)

            # 填充数据
            for row_idx, row in enumerate(rows):
                for col_idx, cell in enumerate(row):
                    if row_idx < num_rows and col_idx < num_cols:
                        cell_text = cell.get("text", "") if isinstance(cell, dict) else str(cell)
                        table.cell(row_idx, col_idx).text = cell_text

                        # 设置样式
                        para = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                        para.font.size = Pt(12)
                        para.font.name = "Microsoft YaHei"

                        # 表头样式
                        if row_idx == 0:
                            para.font.bold = True
                            para.font.color.rgb = RGBColor(255, 255, 255)
                            cell = table.cell(row_idx, col_idx)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(99, 102, 241)

        except Exception:
            pass

    def _add_latex_element(self, slide, element: SlideElement):
        """添加LaTeX公式元素（作为文本显示）"""
        if not element.latex:
            return

        # LaTeX公式作为文本显示（PPTX原生不支持LaTeX）
        textbox = slide.shapes.add_textbox(
            self._px(element.left),
            self._px(element.top),
            self._px(element.width),
            self._px(element.height)
        )
        tf = textbox.text_frame
        tf.paragraphs[0].text = element.latex
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.name = "Cambria Math"
        tf.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
        tf.word_wrap = True

    def _add_line_element(self, slide, element: SlideElement):
        """添加线条元素"""
        points = element.points
        if len(points) < 2:
            return

        try:
            start_x, start_y = points[0]
            end_x, end_y = points[1]

            connector = slide.shapes.add_connector(
                1,  # 直线连接器
                self._px(start_x),
                self._px(start_y),
                self._px(end_x),
                self._px(end_y)
            )

            connector.line.color.rgb = self._parse_color(element.line_color or "#333333")

            # 线条样式
            line_style = element.line_style.lower() if element.line_style else "solid"
            if line_style == "dashed":
                connector.line.dash_style = 2  # DASH
            elif line_style == "dotted":
                connector.line.dash_style = 3  # DOT

            connector.line.width = Pt(2)

        except Exception:
            pass


# 全局实例
_pptx_exporter: Optional[PPTXExporter] = None


def get_pptx_exporter() -> PPTXExporter:
    """获取PPTX导出器实例"""
    global _pptx_exporter
    if _pptx_exporter is None:
        _pptx_exporter = PPTXExporter()
    return _pptx_exporter


def export_course_to_pptx(course_data: CourseData) -> bytes:
    """快捷函数：将课程数据导出为PPTX"""
    exporter = get_pptx_exporter()
    return exporter.export(course_data)


def export_course_to_file(course_data: CourseData, output_path: str) -> str:
    """快捷函数：将课程数据导出到文件"""
    exporter = get_pptx_exporter()
    return exporter.export_to_file(course_data, output_path)